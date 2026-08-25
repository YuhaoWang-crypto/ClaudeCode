#!/usr/bin/env python3
"""
Geometry lint for report.pptx.

LibreOffice cannot load .pptx in this environment, so slide images are not
available for visual QA. This checks the defects that a render would have
caught, directly against the shape geometry:

  * shapes outside the slide
  * text that needs more height than its box has (estimated from an average
    glyph width of 0.50 em and 1.30 line spacing - deliberately generous, so a
    flag means a real risk of clipping rather than a rounding artefact)
  * overlapping text boxes

It is an approximation. Run it, fix what it flags, and treat a clean run as
"no obvious layout defect", not as proof the deck renders perfectly.
"""
import sys

from pptx import Presentation
from pptx.util import Emu

SLIDE_W, SLIDE_H = 13.333, 7.5
EM_WIDTH = 0.50          # average glyph width as a fraction of font size
LINE = 1.30              # line height multiple
TOL = 1.06               # allow 6% over before flagging


def inches(v):
    return Emu(v).inches if v is not None else 0.0


def needed_height(shape):
    """Estimated height in inches for the text at its declared sizes."""
    tf = shape.text_frame
    w = inches(shape.width) - 0.16          # text-box internal padding
    if w <= 0:
        return 0.0
    total = 0.0
    for para in tf.paragraphs:
        text = "".join(r.text for r in para.runs)
        size = max((r.font.size.pt for r in para.runs if r.font.size), default=None)
        if size is None:
            size = para.font.size.pt if para.font.size else 12.0
        if not text:
            total += size * LINE / 72.0
            continue
        cpl = max(int(w * 72.0 / (size * EM_WIDTH)), 1)
        lines = 0
        for hard in text.split("\n"):
            lines += max(-(-len(hard) // cpl), 1)
        total += lines * size * LINE / 72.0
    return total


def rects_overlap(a, b):
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    ox = max(0.0, min(ax + aw, bx + bw) - max(ax, bx))
    oy = max(0.0, min(ay + ah, by + bh) - max(ay, by))
    area = ox * oy
    if area <= 0:
        return 0.0
    return area / max(min(aw * ah, bw * bh), 1e-9)


def main(path="report.pptx"):
    pres = Presentation(path)
    problems = []
    for i, slide in enumerate(pres.slides, 1):
        boxes = []
        for shape in slide.shapes:
            x, y = inches(shape.left), inches(shape.top)
            w, h = inches(shape.width), inches(shape.height)
            name = shape.shape_type
            label = ""
            if shape.has_text_frame and shape.text_frame.text.strip():
                label = shape.text_frame.text.strip().replace("\n", " ")[:44]

            if x < -0.01 or y < -0.01 or x + w > SLIDE_W + 0.01 or y + h > SLIDE_H + 0.01:
                problems.append(f"slide {i}: off-slide {name} at "
                                f"({x:.2f},{y:.2f}) {w:.2f}x{h:.2f} '{label}'")
            if shape.has_text_frame and label:
                need = needed_height(shape)
                if need > h * TOL:
                    problems.append(f"slide {i}: text may overflow - needs "
                                    f"{need:.2f}in, box {h:.2f}in '{label}'")
                boxes.append(((x, y, w, h), label))

        for j in range(len(boxes)):
            for k in range(j + 1, len(boxes)):
                frac = rects_overlap(boxes[j][0], boxes[k][0])
                if frac > 0.20:
                    problems.append(f"slide {i}: text boxes overlap "
                                    f"{frac*100:.0f}% - '{boxes[j][1]}' / '{boxes[k][1]}'")

    if problems:
        print(f"{len(problems)} potential layout issues:")
        for p in problems:
            print("  " + p)
        return 1
    print(f"{len(pres.slides.__iter__.__self__._sldIdLst)} slides, no layout issues found")
    return 0


if __name__ == "__main__":
    sys.exit(main(*sys.argv[1:]))
